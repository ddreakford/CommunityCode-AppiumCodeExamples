#!/usr/bin/env python3
"""
Create PowerPoint presentation with Appium Samples In a Box diagrams
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_powerpoint_presentation():
    """Create PowerPoint presentation with architectural and sequence diagrams"""
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Appium Samples In a Box"
    subtitle.text = "Architecture & Sequence Flow Diagrams"
    
    # Format title
    title_font = title.text_frame.paragraphs[0].font
    title_font.size = Pt(44)
    title_font.bold = True
    title_font.color.rgb = RGBColor(25, 118, 210)  # Blue
    
    # Format subtitle
    subtitle_font = subtitle.text_frame.paragraphs[0].font
    subtitle_font.size = Pt(20)
    subtitle_font.color.rgb = RGBColor(66, 66, 66)  # Dark gray
    
    # Slide 2: System Architecture
    blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add title to architecture slide
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Architectural Overview"
    title_para.alignment = PP_ALIGN.CENTER
    title_font = title_para.font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = RGBColor(25, 118, 210)
    
    # Add architecture diagram
    arch_image_path = 'appium_architecture.png'
    if os.path.exists(arch_image_path):
        slide2.shapes.add_picture(arch_image_path, Inches(0.5), Inches(1), Inches(9), Inches(6))
    
    # Add architecture description
    desc_box = slide2.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_para = desc_frame.paragraphs[0]
    desc_para.text = "Key Components: Runner (Python) • Docker Container • Build Tools (Gradle/uv) • Digital.ai Testing Cloud"
    desc_para.alignment = PP_ALIGN.CENTER
    desc_font = desc_para.font
    desc_font.size = Pt(14)
    desc_font.color.rgb = RGBColor(66, 66, 66)
    
    # Slide 3: Sequence Flow
    slide3 = prs.slides.add_slide(blank_slide_layout)
    
    # Add title to sequence slide
    seq_title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    seq_title_frame = seq_title_box.text_frame
    seq_title_para = seq_title_frame.paragraphs[0]
    seq_title_para.text = "Test Execution Sequence Flow"
    seq_title_para.alignment = PP_ALIGN.CENTER
    seq_title_font = seq_title_para.font
    seq_title_font.size = Pt(32)
    seq_title_font.bold = True
    seq_title_font.color.rgb = RGBColor(25, 118, 210)
    
    # Add sequence diagram
    seq_image_path = 'appium_sequence.png'
    if os.path.exists(seq_image_path):
        slide3.shapes.add_picture(seq_image_path, Inches(0.5), Inches(1), Inches(9), Inches(6))
    
    # Add sequence description
    seq_desc_box = slide3.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.8))
    seq_desc_frame = seq_desc_box.text_frame
    seq_desc_para = seq_desc_frame.paragraphs[0]
    seq_desc_para.text = "Execution Flow: User Command → Environment Validation → Docker → Build Tools → Cloud Connection → Parallel Testing → Reports"
    seq_desc_para.alignment = PP_ALIGN.CENTER
    seq_desc_font = seq_desc_para.font
    seq_desc_font.size = Pt(14)
    seq_desc_font.color.rgb = RGBColor(66, 66, 66)
    
    # Slide 4: Key Features & Configuration
    slide4 = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    features_title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    features_title_frame = features_title_box.text_frame
    features_title_para = features_title_frame.paragraphs[0]
    features_title_para.text = "Key Features & Configuration"
    features_title_para.alignment = PP_ALIGN.CENTER
    features_title_font = features_title_para.font
    features_title_font.size = Pt(32)
    features_title_font.bold = True
    features_title_font.color.rgb = RGBColor(25, 118, 210)
    
    # Add features content
    features_content = """🎯 Test Execution Features:
• Parallel execution with configurable workers (default: 4)
• Support for Java/TestNG and Python/pytest test suites
• Docker containerized execution environment
• HTML and JSON report generation
• Real-time test output streaming

🔧 Environment Configuration:
• CLOUD_URL: https://uscloud.experitest.com
• ACCESS_KEY: Digital.ai Testing Cloud access key
• APPIUM_VERSION: 2.0.0 (configurable)
• ANDROID_DEVICE_QUERY: Device selection filters
• IOS_DEVICE_QUERY: Device selection filters

🐳 Docker Integration:
• Volume mounting for reports and logs persistence
• Environment variable injection from .env file
• Resource limits and parallel execution control
• Separate dev and runtime container profiles

📊 Report Generation:
• HTML test reports (Java & Python)
• JSON summary with execution statistics
• Test logs with timestamps
• Mounted volumes for external access"""
    
    content_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    content_frame = content_box.text_frame
    content_para = content_frame.paragraphs[0]
    content_para.text = features_content
    content_font = content_para.font
    content_font.size = Pt(14)
    content_font.color.rgb = RGBColor(66, 66, 66)
    
    # Save presentation
    pptx_path = 'Appium_Samples_InaBox_Architecture.pptx'
    prs.save(pptx_path)
    
    return pptx_path

def main():
    """Generate PowerPoint presentation"""
    print("Creating PowerPoint presentation...")
    
    pptx_file = create_powerpoint_presentation()
    
    print(f"✅ PowerPoint presentation created: {os.path.basename(pptx_file)}")
    print(f"📁 Full path: {pptx_file}")
    print("\n📋 Presentation includes:")
    print("  1. Title slide with system overview")
    print("  2. System Architecture diagram")
    print("  3. Test Execution Sequence flow")
    print("  4. Key Features & Configuration details")

if __name__ == "__main__":
    main()